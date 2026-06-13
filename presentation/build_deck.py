"""
Builds the SecureAIExam hackathon deck (16:9) as a single .pptx.

Design language mirrors the app: deep-navy canvas, blue→violet→cyan accents,
gradient accent bars, glass-ish cards. Re-run any time to regenerate:

    python presentation/build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (from mobile/src/lib/theme.ts dark mode) -----------------------
BG       = RGBColor(0x06, 0x0B, 0x16)
PANEL    = RGBColor(0x0E, 0x17, 0x29)
PANEL2   = RGBColor(0x15, 0x21, 0x3A)
BORDER   = RGBColor(0x2C, 0x3E, 0x63)
TEXT     = RGBColor(0xED, 0xF2, 0xFC)
DIM      = RGBColor(0x9D, 0xB0, 0xCE)
FAINT    = RGBColor(0x6C, 0x7E, 0xA1)
PRIMARY  = RGBColor(0x5B, 0x8D, 0xF8)
ACCENT   = RGBColor(0x22, 0xD3, 0xEE)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
SUCCESS  = RGBColor(0x34, 0xD3, 0x99)
WARN     = RGBColor(0xFB, 0xBF, 0x24)
DANGER   = RGBColor(0xF8, 0x71, 0x71)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    _no_line(sp)
    return sp


def gradient(s, x, y, w, h, c1, c2, angle=0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    _no_line(sp)
    try:
        sp.fill.gradient()
        stops = sp.fill.gradient_stops
        stops[0].color.rgb = c1
        stops[0].position = 0.0
        stops[1].color.rgb = c2
        stops[1].position = 1.0
        try:
            sp.fill.gradient_angle = angle
        except Exception:
            pass
    except Exception:
        sp.fill.solid()
        sp.fill.fore_color.rgb = c1
    return sp


def tri_gradient(s, x, y, w, h):
    """Brand bar: cyan -> blue -> violet."""
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _no_line(sp)
    sp.fill.gradient()
    g = sp.fill._xPr.find(qn('a:gradFill'))
    lst = g.find(qn('a:gsLst'))
    for gs in list(lst):
        lst.remove(gs)
    for pos, col in [(0, ACCENT), (50000, PRIMARY), (100000, VIOLET)]:
        gs = lst.makeelement(qn('a:gs'), {'pos': str(pos)})
        srgb = gs.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(srgb)
        lst.append(gs)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line
        for txt, size, color, bold in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Segoe UI"
    return tb


def kicker(s, label, color=ACCENT):
    rect(s, Inches(0), Inches(0), SW, Inches(0.14), BG)  # spacer (keeps z-order tidy)
    tri_gradient(s, Inches(0.6), Inches(0.62), Inches(0.5), Inches(0.09))
    text(s, Inches(0.6), Inches(0.72), Inches(11), Inches(0.4),
         [[(label.upper(), 13, color, True)]])


def title(s, t, y=Inches(1.06), size=33):
    text(s, Inches(0.6), y, Inches(12.1), Inches(1.0),
         [[(t, size, TEXT, True)]])


def footer(s, n):
    text(s, Inches(0.6), Inches(7.0), Inches(6), Inches(0.4),
         [[("SecureAIExam", 10, FAINT, True), ("  ·  leak-proof examinations", 10, FAINT, False)]])
    text(s, Inches(11.0), Inches(7.0), Inches(1.7), Inches(0.4),
         [[(str(n), 10, FAINT, True)]], align=PP_ALIGN.RIGHT)


def bullets(s, items, x=Inches(0.65), y=Inches(2.0), w=Inches(12.0), size=17,
            gap=12, marker="▸", mcolor=ACCENT):
    tb = s.shapes.add_textbox(x, y, w, Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, sub = item
        else:
            head, sub = item, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        m = p.add_run(); m.text = marker + "  "
        m.font.size = Pt(size); m.font.color.rgb = mcolor; m.font.bold = True; m.font.name = "Segoe UI"
        r = p.add_run(); r.text = head
        r.font.size = Pt(size); r.font.color.rgb = TEXT; r.font.bold = True; r.font.name = "Segoe UI"
        if sub:
            r2 = p.add_run(); r2.text = "  —  " + sub
            r2.font.size = Pt(size); r2.font.color.rgb = DIM; r2.font.bold = False; r2.font.name = "Segoe UI"
    return tb


def card(s, x, y, w, h, accent, head, body, num=None):
    rect(s, x, y, w, h, PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, Inches(0.09), h, accent)  # left accent spine
    pad = Inches(0.28)
    if num:
        text(s, x + pad, y + Inches(0.16), w - pad, Inches(0.5),
             [[(num, 14, accent, True)]])
        hy = y + Inches(0.5)
    else:
        hy = y + Inches(0.22)
    text(s, x + pad, hy, w - pad - Inches(0.1), Inches(0.5),
         [[(head, 16, TEXT, True)]])
    text(s, x + pad, hy + Inches(0.46), w - pad - Inches(0.1), h - Inches(1.0),
         [[(body, 12.5, DIM, False)]], line=1.1)


# ============================================================ SLIDE 1 — TITLE
s = slide()
gradient(s, Inches(0), Inches(0), SW, SH, RGBColor(0x07, 0x0E, 0x20), BG, angle=90)
# decorative aurora blobs
b = rect(s, Inches(8.6), Inches(-1.4), Inches(6), Inches(6), VIOLET, MSO_SHAPE.OVAL); b.fill.fore_color.rgb = RGBColor(0x18, 0x12, 0x3A)
b = rect(s, Inches(-1.6), Inches(3.8), Inches(5.5), Inches(5.5), PRIMARY, MSO_SHAPE.OVAL); b.fill.fore_color.rgb = RGBColor(0x0B, 0x1A, 0x39)
tri_gradient(s, Inches(0.9), Inches(2.05), Inches(1.5), Inches(0.13))
text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5),
     [[("🛡  SECURE  ·  AI  ·  EXAM", 15, ACCENT, True)]])
text(s, Inches(0.86), Inches(2.75), Inches(11.6), Inches(1.5),
     [[("SecureAIExam", 64, WHITE, True)]])
text(s, Inches(0.9), Inches(3.95), Inches(11.0), Inches(0.7),
     [[("Leak-proof examinations at national scale", 26, TEXT, True)]])
text(s, Inches(0.9), Inches(4.7), Inches(11.4), Inches(0.6),
     [[("Cryptographic dual-control", 15, PRIMARY, True),
       ("   ·   ", 15, FAINT, False),
       ("Per-copy forensic identity", 15, ACCENT, True),
       ("   ·   ", 15, FAINT, False),
       ("Real-time chain of custody", 15, VIOLET, True)]])
text(s, Inches(0.9), Inches(6.5), Inches(11.4), Inches(0.5),
     [[("Hackathon 2026", 13, DIM, True),
       ("    github.com/dynagi/Exam-Security", 13, FAINT, False)]])

# ============================================================ SLIDE 2 — PROBLEM
s = slide()
kicker(s, "The Problem", DANGER)
title(s, "Question-paper leaks are a national crisis")
bullets(s, [
    ("NEET, UPSC & Board leaks repeatedly derail exams", "for millions of aspirants each year"),
    ("Leaks happen at every stage", "authoring, printing, transport, storage, the exam hall"),
    ("Today's process trusts people and paper", "every stage is a single point of failure"),
    ("Once a paper is out, it's untraceable", "no way to know which copy leaked or who held it"),
    ("The fallout is brutal", "cancelled exams, eroded public trust, lost years for honest students"),
], mcolor=DANGER)
footer(s, 2)

# ============================================================ SLIDE 3 — PROBLEM STATEMENT
s = slide()
kicker(s, "Problem Statement", WARN)
title(s, "What we set out to solve")
rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.5), PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(2.1), Inches(0.12), Inches(2.5), WARN)
text(s, Inches(1.05), Inches(2.45), Inches(11.4), Inches(1.9),
     [[("Design an end-to-end examination platform that makes a question-paper "
        "leak ", 21, TEXT, False),
       ("cryptographically hard before the exam", 21, ACCENT, True),
       (" and ", 21, TEXT, False),
       ("forensically traceable after it", 21, PRIMARY, True),
       (" — eliminating every single point of failure across authoring, sealing, "
        "printing, transport and the exam center, while proving accountability "
        "with a ", 21, TEXT, False),
       ("tamper-evident record.", 21, VIOLET, True)]], line=1.15)
text(s, Inches(0.65), Inches(5.1), Inches(12), Inches(1.2),
     [[("Success = ", 15, DIM, True),
       ("prevention", 15, SUCCESS, True), ("  +  ", 15, FAINT, False),
       ("detection", 15, ACCENT, True), ("  +  ", 15, FAINT, False),
       ("accountability", 15, VIOLET, True),
       ("   — engineered together, not bolted on.", 15, DIM, False)]])
footer(s, 3)

# ============================================================ SLIDE 4 — WHY EXISTING FIXES FAIL
s = slide()
kicker(s, "Why partial fixes fail")
title(s, "Point solutions leave the door open")
bullets(s, [
    ("Encryption alone", "one admin still holds the key — that's one point of failure"),
    ("Watermarks alone", "they identify a leak, but never prevent early access"),
    ("Manual custody registers", "forgeable, slow, and blind to real-time theft"),
    ("Audit logs you can edit", "useless the moment an insider rewrites history"),
])
rect(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.95), PANEL2, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(1.0), Inches(5.72), Inches(11.4), Inches(0.7),
     [[("The gap: ", 16, ACCENT, True),
       ("no system combines dual-control prevention, per-copy forensics, "
        "live custody and tamper-proof accountability — at once.", 16, TEXT, False)]])
footer(s, 4)

# ============================================================ SLIDE 5 — SOLUTION (4 pillars)
s = slide()
kicker(s, "Our Solution", SUCCESS)
title(s, "Remove every single point of failure")
cy, ch = Inches(2.15), Inches(2.15)
cw, gapx = Inches(2.92), Inches(0.16)
xs = Inches(0.6)
pillars = [
    (PRIMARY, "Dual-control crypto", "AES-256-GCM seal; key XOR-split between Setter & Admin. No one person can open a paper."),
    (ACCENT, "Per-copy identity", "Every printed copy carries a signed QR + an invisible forensic fingerprint."),
    (VIOLET, "Live chain of custody", "Every handoff is scanned and streamed to admins over WebSockets in <1s."),
    (SUCCESS, "Tamper-evident audit", "Hash-chained, append-only history — any edit breaks every later hash."),
]
for i, (acc, h, body) in enumerate(pillars):
    x = Inches(0.6 + i * (2.92 + 0.16))
    card(s, x, cy, cw, ch, acc, h, body)
rect(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.4), BG)
text(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.1),
     [[("Together they make a leak ", 17, DIM, False),
       ("hard to commit", 17, SUCCESS, True),
       (", ", 17, DIM, False),
       ("impossible to hide", 17, ACCENT, True),
       (", and ", 17, DIM, False),
       ("trivial to trace.", 17, PRIMARY, True)]], align=PP_ALIGN.CENTER)
footer(s, 5)

# ============================================================ SLIDE 6 — KEY INNOVATIONS
s = slide()
kicker(s, "Key Innovations")
title(s, "The mechanisms behind leak-proofing")
bullets(s, [
    ("Split-key sealing", "K = random AES-256 key, XOR-split; Setter + Admin passphrases both required"),
    ("Signed + fingerprinted copies", "HMAC-signed QR (forgeries rejected) + steganographic per-copy watermark"),
    ("Real-time alarms", "missing/unscanned copy → CRITICAL alert on every admin device in <1 second"),
    ("Pre-exam reconciliation sweep", "unscanned copies auto-flagged ~20 min before the slot starts"),
    ("Leak forensics from a photo", "AI extracts the fingerprint → exact copy → full custody trail → exact holder"),
    ("Hybrid AI PDF intake", "Claude NLP imports questions from PDFs, with an offline parser fallback"),
], size=16, gap=11)
footer(s, 6)

# ============================================================ SLIDE 7 — LIFECYCLE
s = slide()
kicker(s, "Lifecycle")
title(s, "From question to exam hall — secured at every step")
steps = [
    ("1", PRIMARY, "Author", "Teachers upload Qs (typed or AI PDF import); setter builds the paper"),
    ("2", PRIMARY, "Seal", "Setter seals the paper with passphrase — share 1 of the key"),
    ("3", ACCENT, "Co-sign", "Admin co-signs — share 2. Now DUAL-SEALED: nobody can open it alone"),
    ("4", ACCENT, "Print ceremony", "Both passphrases → per-copy QR + invisible fingerprint minted"),
    ("5", VIOLET, "Custody", "Allocate to centers; invigilator scans each QR — admin notified live"),
    ("6", WARN, "Reconcile", "Unscanned copies → CRITICAL alert ~20 min before the slot"),
    ("7", DANGER, "Forensics", "A leaked photo is traced to the exact copy and its last holder"),
]
y = Inches(1.95)
for n, acc, head, body in steps:
    rect(s, Inches(0.62), y + Inches(0.02), Inches(0.34), Inches(0.34), acc, MSO_SHAPE.OVAL)
    text(s, Inches(0.62), y - Inches(0.04), Inches(0.34), Inches(0.4),
         [[(n, 13, BG, True)]], align=PP_ALIGN.CENTER)
    text(s, Inches(1.15), y - Inches(0.04), Inches(2.3), Inches(0.5),
         [[(head, 15.5, TEXT, True)]])
    text(s, Inches(3.5), y - Inches(0.02), Inches(9.2), Inches(0.5),
         [[(body, 14, DIM, False)]])
    y = y + Inches(0.68)
footer(s, 7)

# ============================================================ SLIDE 8 — ARCHITECTURE
s = slide()
kicker(s, "Architecture")
title(s, "Three trust boundaries, one realtime backbone")
comps = [
    (PRIMARY, "Mobile — Expo / React Native", "Teacher, Setter, Admin & Invigilator UIs from one codebase (iOS · Android · Web)"),
    (ACCENT, "Supabase", "Auth + JWT · Postgres with Row-Level Security · Realtime (WAL → WebSocket fan-out)"),
    (VIOLET, "Node crypto service", "Sealing, co-sign & print ceremonies — service-role only; secrets never touch a client"),
    (SUCCESS, "Python AI service (FastAPI)", "Steganographic fingerprinting · custody anomaly scoring · PDF question NLP"),
]
y = Inches(2.0)
for acc, head, body in comps:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(0.95), PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), y, Inches(0.1), Inches(0.95), acc)
    text(s, Inches(0.95), y + Inches(0.13), Inches(11.5), Inches(0.4),
         [[(head, 16, TEXT, True)]])
    text(s, Inches(0.95), y + Inches(0.5), Inches(11.5), Inches(0.4),
         [[(body, 13, DIM, False)]])
    y = y + Inches(1.08)
text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
     [[("Why split? ", 13.5, ACCENT, True),
       ("secrets stay server-side, the app is untrusted, and AI scales independently on exam day.",
        13.5, DIM, False)]])
footer(s, 8)

# ============================================================ SLIDE 9 — THREAT -> MITIGATION
s = slide()
kicker(s, "Security Model")
title(s, "Threat → Mitigation")
rows = [
    ("Insider opens the paper early", "Dual-control key split — no single passphrase decrypts anything", SUCCESS),
    ("Database breach", "Only ciphertext + passphrase-wrapped shares; RLS hides secrets from every JWT", SUCCESS),
    ("Copy stolen in transit", "Per-copy custody chain — a missed checkpoint = instant alert + last holder", SUCCESS),
    ("Leaked photo on social media", "Invisible fingerprint → exact copy → exact custody trail → exact person", SUCCESS),
    ("Forged 'official' copies", "QR payloads are HMAC-signed; verification fails for fakes", SUCCESS),
    ("Log tampering to hide tracks", "Hash-chained, append-only audit log; any edit breaks every later hash", SUCCESS),
]
y = Inches(1.95)
for threat, mit, acc in rows:
    text(s, Inches(0.65), y, Inches(4.5), Inches(0.6),
         [[("✕  ", 14, DANGER, True), (threat, 14, TEXT, True)]])
    text(s, Inches(5.4), y, Inches(7.3), Inches(0.6),
         [[("→  ", 14, acc, True), (mit, 13.5, DIM, False)]])
    y = y + Inches(0.74)
footer(s, 9)

# ============================================================ SLIDE 10 — TECH STACK
s = slide()
kicker(s, "Tech Stack")
title(s, "Built to ship — and to scale")
groups = [
    (PRIMARY, "Frontend", "React Native · Expo SDK 56 · one codebase → iOS, Android & Web · premium glass UI, animated gradients, light/dark"),
    (ACCENT, "Data & Realtime", "Supabase — Postgres, Row-Level Security, Realtime WebSockets, JWT auth"),
    (VIOLET, "Cryptography", "Node.js · AES-256-GCM · scrypt key-wrap · HMAC-SHA256 · XOR key-split · hash-chained logs"),
    (SUCCESS, "AI / Forensics", "Python FastAPI · steganographic watermarking · custody anomaly scoring · Claude NLP (hybrid + offline)"),
]
y = Inches(2.05)
for acc, head, body in groups:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), y, Inches(0.1), Inches(1.0), acc)
    text(s, Inches(0.95), y + Inches(0.16), Inches(3.0), Inches(0.6),
         [[(head, 15.5, acc, True)]])
    text(s, Inches(3.4), y + Inches(0.16), Inches(9.1), Inches(0.7),
         [[(body, 13.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(1.14)
footer(s, 10)

# ============================================================ SLIDE 11 — DEMO
s = slide()
kicker(s, "Live Demo", DANGER)
title(s, "The money shots")
bullets(s, [
    ("Dual-seal, live", "watch a paper flip to DUAL-SEALED across two devices the instant the admin co-signs"),
    ("Scan-in, live", "scan a copy's QR at the 'center' → admin dashboard pings instantly"),
    ("Reconciliation alarm", "leave a copy unscanned → CRITICAL alert fires before the slot"),
    ("Report MISSING", "one tap → red alert on every admin device in <1 second"),
    ("Open the audit log", "the hash-chained, tamper-evident story of everything that happened"),
], mcolor=DANGER, size=16, gap=12)
footer(s, 11)

# ============================================================ SLIDE 12 — IMPACT
s = slide()
kicker(s, "Impact & Scale", SUCCESS)
title(s, "Why it matters — and that it holds up")
left = [
    ("Protects millions of aspirants", "from cancelled, compromised, re-scheduled exams"),
    ("Restores trust in the system", "leaks become provable, attributable, rare"),
    ("Runs on existing devices", "invigilators just use a phone camera"),
]
right = [
    ("Custody scans are single-row inserts", "tens of thousands/min; partition by exam"),
    ("Stateless AI & crypto services", "horizontal scale behind a load balancer on exam day"),
    ("Realtime fan-out is filtered per-JWT", "admins subscribe only to their channels"),
]
text(s, Inches(0.6), Inches(1.95), Inches(6), Inches(0.4), [[("HUMAN IMPACT", 13, SUCCESS, True)]])
bullets(s, left, x=Inches(0.6), y=Inches(2.4), w=Inches(5.9), size=14.5, gap=10, mcolor=SUCCESS)
text(s, Inches(6.9), Inches(1.95), Inches(6), Inches(0.4), [[("ENGINEERING SCALE", 13, ACCENT, True)]])
bullets(s, right, x=Inches(6.9), y=Inches(2.4), w=Inches(5.9), size=14.5, gap=10, mcolor=ACCENT)
footer(s, 12)

# ============================================================ SLIDE 13 — ROADMAP
s = slide()
kicker(s, "Production Roadmap")
title(s, "Honest about demo trade-offs → production")
rows = [
    ("Escrow window between seal & co-sign", "shareB → Admin hardware token / true Shamir t-of-n with an HSM ceremony"),
    ("LSB watermark (demo)", "robust DCT/DWT watermarking that survives print → photograph"),
    ("Self-selected roles (demo)", "invite-only provisioning by the examination board"),
    ("Plaintext question bank", "purge or re-encrypt the bank after sealing"),
]
y = Inches(2.1)
for now, nxt in rows:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.95), y + Inches(0.16), Inches(5.2), Inches(0.7),
         [[("Now  ", 12, WARN, True), (now, 14, TEXT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(6.2), y + Inches(0.3), Inches(0.6), Inches(0.4), [[("→", 18, ACCENT, True)]])
    text(s, Inches(6.9), y + Inches(0.16), Inches(5.6), Inches(0.7),
         [[("Prod  ", 12, SUCCESS, True), (nxt, 13.5, DIM, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(1.14)
footer(s, 13)

# ============================================================ SLIDE 14 — CLOSE
s = slide()
gradient(s, Inches(0), Inches(0), SW, SH, RGBColor(0x07, 0x0E, 0x20), BG, angle=90)
b = rect(s, Inches(8.8), Inches(2.2), Inches(6), Inches(6), VIOLET, MSO_SHAPE.OVAL); b.fill.fore_color.rgb = RGBColor(0x16, 0x12, 0x36)
tri_gradient(s, Inches(0.9), Inches(2.5), Inches(1.5), Inches(0.13))
text(s, Inches(0.86), Inches(2.75), Inches(11.6), Inches(1.2),
     [[("Make leaks ", 40, WHITE, True),
       ("cryptographically hard", 40, ACCENT, True)]])
text(s, Inches(0.86), Inches(3.6), Inches(11.6), Inches(1.2),
     [[("— and ", 40, WHITE, True),
       ("forensically traceable.", 40, PRIMARY, True)]])
text(s, Inches(0.9), Inches(4.9), Inches(11), Inches(0.6),
     [[("SecureAIExam", 20, TEXT, True),
       ("  ·  dual-control · per-copy forensics · live custody · tamper-proof audit", 16, DIM, False)]])
text(s, Inches(0.9), Inches(6.2), Inches(11), Inches(0.6),
     [[("Thank you  —  Questions?", 18, VIOLET, True)]])
text(s, Inches(0.9), Inches(6.8), Inches(11), Inches(0.4),
     [[("github.com/dynagi/Exam-Security", 12, FAINT, False)]])

import os
out = os.path.join(os.path.dirname(__file__), "SecureAIExam_Hackathon.pptx")
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slides")
